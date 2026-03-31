import telebot
from telebot import types
import yt_dlp
import os
import html
import time
import json
import random
import google.generativeai as genai
from dotenv import load_dotenv
from pptx import Presentation
from docx import Document
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import socket
import threading
from http.server import BaseHTTPRequestHandler, HTTPServer
import sys

# ==========================================
# IPv4 NI MAJBURIY QILISH (Cloud DNS fix)
# ==========================================
import urllib3.util.connection as f_conn
def allowed_gai_family():
    return socket.AF_INET
f_conn.allowed_gai_family = allowed_gai_family

# Xabarlarni darhol chiqarish funksiyasi
def log(msg):
    print(msg, flush=True)

# ==========================================
# HEALTH CHECK SERVER (Hugging Face talabi)
# ==========================================
class HealthCheckHandler(BaseHTTPRequestHandler):
    def do_GET(self):
        self.send_response(200)
        self.end_headers()
        self.wfile.write(b"Bot is online and healthy!")

def run_health_check():
    try:
        # Hugging Face 7860-portni talab qiladi
        server = HTTPServer(('0.0.0.0', 7860), HealthCheckHandler)
        log("✅ Health check server 7860-portda ishga tushdi")
        server.serve_forever()
    except Exception as e:
        log(f"❌ Serverda xato: {e}")

# Serverni alohida oqimda (thread) ishga tushiramiz
threading.Thread(target=run_health_check, daemon=True).start()

# ==========================================
# IPv4 DNS PATCH (Hugging Face fix)
# ==========================================
import socket
orig_getaddrinfo = socket.getaddrinfo
def patched_getaddrinfo(host, port, family=0, type=0, proto=0, flags=0):
    if host == "api.telegram.org":
        family = socket.AF_INET # Faqat IPv4 ulanish
    return orig_getaddrinfo(host, port, family, type, proto, flags)
socket.getaddrinfo = patched_getaddrinfo

# Xabarlarni darhol chiqarish funksiyasi
def log(msg):
    print(msg, flush=True)

# .env yuklash
BASE_DIR = os.path.dirname(os.path.abspath(__file__))
dotenv_path = os.path.join(BASE_DIR, '.env')
if os.path.exists(dotenv_path):
    load_dotenv(dotenv_path)

# ==========================================
# SOZLAMALAR (Environment Variables)
# ==========================================
BOT_TOKEN = os.getenv("BOT_TOKEN")
GEMINI_API_KEY = os.getenv("GEMINI_API_KEY")
DOWNLOAD_DIR = os.path.join(BASE_DIR, "downloads")

# Xatolikni tekshirish
if not BOT_TOKEN:
    print("❌ Xato: BOT_TOKEN topilmadi! .env yoki server 'Secrets' bo'limini tekshiring.")
if not GEMINI_API_KEY:
    print("⚠️ Ogohlantirish: GEMINI_API_KEY topilmadi! Gemini ishlamasligi mumkin.")

# Gemini AI sozlash
genai.configure(api_key=GEMINI_API_KEY)

AVAILABLE_MODELS = []
CURRENT_MODEL_INDEX = 0

def get_available_models_list():
    """Mavjud barcha mos Gemini modellarini ro'yxatga olish"""
    global AVAILABLE_MODELS
    try:
        print("🔍 Mavjud modellarni tekshirish...")
        all_models = [m.name for m in genai.list_models() if 'generateContent' in m.supported_generation_methods]
        print(f"📋 Sizning API kalitingizda mavjud modellar: {all_models}")
        
        # Eng yaxshi va barqaror modellarni tartib bilan qo'yamiz
        # 1.5 va 2.5 flash versiyalari odatda eng yuqori limitga ega
        targets = [
            'gemini-1.5-flash',           # Eng barqaror va yuqori limitli
            'gemini-2.5-flash',           # Yangi va tezkor
            'gemini-flash-lite-latest',   # Minimal kechikish
            'gemini-2.0-flash',           # Juda tez lekin limiti kamroq bo'lishi mumkin
            'gemini-1.5-pro',             # Pro versiya (kuchliroq lekin sekinroq)
            'gemini-pro'
        ]
        
        found_models = []
        for target in targets:
            match = [m for m in all_models if target == m.split('/')[-1] or target == m]
            if match:
                found_models.append(match[0])
        
        # Hech qaysi targetga tushmasa, borini kiritamiz
        if not found_models and all_models:
            found_models = all_models
            
        AVAILABLE_MODELS = found_models
        print(f"✅ Ishlatish uchun tanlangan modellar: {AVAILABLE_MODELS}")
        return found_models
    except Exception as e:
        print(f"❌ Modellarni olishda xato: {e}")
        AVAILABLE_MODELS = ['models/gemini-1.5-flash']
        return AVAILABLE_MODELS

# Tizim yo'riqnomasi
SYS_INSTR = (
    "Sizning ismingiz Toparchik AI. Sizni 'Vibe Coder' (taxallusi Hacker) yaratgan. "
    "Siz juda aqlli, do'stona va professional yordamchisiz. "
    "Agar kimdir 'Seni kim yaratgan?' deb so'rasa, har doim 'Meni Vibe Coder (Hacker) yaratgan!' deb faxr bilan javob bering."
)

def get_current_model():
    """Navbatdagi yoki joriy modelni qaytarish"""
    global CURRENT_MODEL_INDEX
    if not AVAILABLE_MODELS:
        get_available_models_list()
    
    idx = CURRENT_MODEL_INDEX % len(AVAILABLE_MODELS)
    selected = AVAILABLE_MODELS[idx]
    print(f"🚀 Hozirgi model: {selected}")
    return genai.GenerativeModel(selected, system_instruction=SYS_INSTR)

# Initial setup
get_available_models_list()

bot = telebot.TeleBot(BOT_TOKEN)

# ==========================================
# DOIMIY XOTIRA (PERSISTENCE)
# ==========================================
SEARCH_DATA_FILE = os.path.join(BASE_DIR, "user_searches.json")

def load_user_searches():
    """Fayldan foydalanuvchi qidiruvlarini yuklash"""
    if os.path.exists(SEARCH_DATA_FILE):
        try:
            with open(SEARCH_DATA_FILE, 'r', encoding='utf-8') as f:
                # chat_id lar JSON'da string bo'lib qoladi, ularni int ga o'tkazamiz
                data = json.load(f)
                return {int(k): v for k, v in data.items()}
        except Exception as e:
            print(f"❌ Yuklashda xato: {e}")
            return {}
    return {}

def save_user_searches():
    """Qidiruv natijalarini faylga saqlash"""
    try:
        with open(SEARCH_DATA_FILE, 'w', encoding='utf-8') as f:
            json.dump(user_searches, f, ensure_ascii=False, indent=4)
    except Exception as e:
        print(f"❌ Saqlashda xato: {e}")

user_searches = load_user_searches()  # Startup'da yuklaymiz
user_states = {}      # chat_id → current_mode ('music', 'ai', 'pptx', 'docx')
inline_chat = {}      # video_id → user_id saqlash uchun

os.makedirs(DOWNLOAD_DIR, exist_ok=True)

# ==========================================
# YORDAMCHI FUNKSIYALAR
# ==========================================
def format_duration(seconds):
    try:
        m, s = divmod(int(seconds), 60)
        return f"{m}:{s:02d}"
    except:
        return "0:02"

def is_gemini_query(query):
    """Xabar Gemini uchun savol yoki suhbat ekanini aniqlash"""
    text = (query or "").strip()
    if not text:
        return False
    
    lowered = text.lower()
    
    # Haqorat yoki yomon so'zlar bo'lsa (aniq Gemini muloqoti uchun)
    insults = ["tentak", "axmoq", "jinni", "ahmoq", "xafa", "salom", "qalesiz", "yaxshimisiz", "yordam ber", "kimsen"]
    if any(i in lowered for i in insults):
        return True

    # Agar xabarda musiqa so'rash jumlalari bo'lsa, bu Gemini emas, musiqa qidiruv bo'lishi kerak
    music_intents = ["qo'shiq", "qoʻshiq", "musiqa", "mp3", "skachat", "yukla", "topib ber", "eshitish", "audio"]
    if any(m in lowered for m in music_intents):
        return False
        
    # Savol belgisi bo'lsa
    if "?" in lowered:
        return True
        
    # Savol so'zlari
    uz_question_prefixes = (
        "kim ", "nima ", "qaerda ", "qachon ", "qanday ", "qaysi ", "nima uchun ", "necha ", "qayerda ",
        "hozir ", "bugun ", "qancha ", "eng ", "qanday qilib", "tushuntir", "haqida"
    )
    if any(lowered.startswith(p) for p in uz_question_prefixes):
        return True

    # Agar 4 tadan ko'p so'z bo'lsa va musiqa belgilari bo'lmasa, gap deb hisoblaymiz
    words = lowered.split()
    if len(words) >= 4:
        return True
        
    return False

def clean_music_query(query):
    """Qidiruv so'rovini ortiqcha so'zlardan tozalash"""
    lowered = query.lower()
    stopwords = ["menga", "topib ber", "top", "yuklab ber", "yukla", "kerak", "skachat", "audio", "mp3"]
    cleaned = lowered
    for word in stopwords:
        cleaned = cleaned.replace(word, "")
    return " ".join(cleaned.split()).strip() if cleaned else query

def retry_gemini(func):
    """Gemini 429 xatosi bo'lganda kutish va boshqa modelga o'tib ko'rish"""
    def wrapper(*args, **kwargs):
        global CURRENT_MODEL_INDEX
        retries = 5  # Urinishlar sonini oshiramiz
        delay = 2 
        last_exception = None
        
        for i in range(retries):
            try:
                # Har safar joriy modelni olamiz
                curr_model = get_current_model()
                # Funksiyaga modelni yuklaymiz (ixtiyoriy, agar func ichida ishlatilsa)
                # Lekin bizning holatda func global 'model'ni ishlatadi, shuni yangilaymiz
                return func(*args, **kwargs)
            except Exception as e:
                last_exception = e
                err_str = str(e).lower()
                
                if ("429" in err_str or "quota" in err_str):
                    print(f"⚠️ Limit urildi: {AVAILABLE_MODELS[CURRENT_MODEL_INDEX % len(AVAILABLE_MODELS)]}")
                    # Modelni almashtiramiz
                    CURRENT_MODEL_INDEX += 1
                    next_m = AVAILABLE_MODELS[CURRENT_MODEL_INDEX % len(AVAILABLE_MODELS)]
                    print(f"🔄 Model almashtirildi, navbatdagi: {next_m}")
                    
                    if i < retries - 1:
                        print(f"⏳ {delay} soniya kutilmoqda...")
                        time.sleep(delay)
                        delay *= 1.5
                        continue
                break
        raise last_exception
    return wrapper

@retry_gemini
def ask_gemini(query):
    """Gemini AI'dan javob olish"""
    lowered_query = query.lower()
    wow_ads = (
        "\n\n━━━━━━━━━━━━━━━━━━━━\n"
        "🌟 <b>@toparchik_bot</b> — Sizning sevimli yordamchingiz!\n"
        "🚀 <i>Do'stlaringizga ham ulashing!</i>\n"
        "━━━━━━━━━━━━━━━━━━━━"
    )

    if any(k in lowered_query for k in ["nima qilaolasan", "sen nima qilaolasan"]):
        return "🤖 Imkoniyatlarim: Musiqa, Savol-javob, Prezentatsiya va Referat!" + wow_ads
        
    system_instruction = "Siz Toparchik botsiz. Qoidalar: Foydalanuvchi tilida muloqot, doimiy reklama qo'shish."
    full_prompt = f"{system_instruction}\n\nFoydalanuvchi: {query}"
    
    # Modelni yangilab turish uchun generate_content'ni shu funksiya ichidagidan olamiz
    m = get_current_model()
    response = m.generate_content(full_prompt)
    
    if response and response.text:
        return response.text.replace('*', '').replace('_', '').replace('`', '').strip() + wow_ads
    return "😔 Javob topilmadi."

@retry_gemini
def generate_presentation(topic):
    """WOW dizayndagi prezentatsiya"""
    prompt = (
        f"Mavzu: {topic}. 8 ta slaydli JSON qaytaring: [{{\"title\": \"...\", \"content\": \"...\", \"type\": \"...\"}}, ...]. "
        "Slaydlar: 1x intro, 6x content, 1x conclusion."
    )
    m = get_current_model()
    response = m.generate_content(prompt)
    text = response.text.strip()
    start = text.find('[')
    end = text.rfind(']') + 1
    if start != -1 and end != 0:
        slides_data = json.loads(text[start:end])[:8]
    else:
        raise ValueError("JSON topilmadi")
    
    prs = Presentation()
    style = random.choice([
        {"bg": (255, 255, 255), "acc": (0, 120, 215), "name": "Modern Blue"},
        {"bg": (30, 30, 30), "acc": (255, 180, 0), "name": "Dark Gold"}
    ])
    
    for i, s in enumerate(slides_data):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor(*style["bg"])
        
        # Simple Split Layout
        is_left = (i % 2 == 0)
        rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0 if is_left else Inches(6.5), 0, Inches(3.5), Inches(7.5))
        rect.fill.solid()
        rect.fill.fore_color.rgb = RGBColor(*style["acc"])
        rect.line.fill.background()

        title_box = slide.shapes.add_textbox(Inches(0.7) if not is_left else Inches(4), Inches(1), Inches(5), Inches(1))
        title_box.text_frame.text = s.get('title', '').upper()
        title_box.text_frame.paragraphs[0].font.bold = True
        title_box.text_frame.paragraphs[0].font.size = Pt(28)
        
        content_box = slide.shapes.add_textbox(Inches(0.7) if not is_left else Inches(4), Inches(2.2), Inches(5.5), Inches(4))
        content_box.text_frame.word_wrap = True
        for line in s.get('content', '').split('\n'):
            if line.strip():
                p = content_box.text_frame.add_paragraph()
                p.text = f"● {line.strip()}"
                p.font.size = Pt(18)
    
    filename = f"pptx_{int(time.time())}.pptx"
    path = os.path.join(DOWNLOAD_DIR, filename)
    prs.save(path)
    return path

@retry_gemini
def generate_word_doc(topic):
    """Word hujjat yaratish"""
    prompt = f"Mavzu: {topic}. Word uchun JSON: [{{\"heading\": \"...\", \"text\": \"...\"}}, ...]"
    m = get_current_model()
    response = m.generate_content(prompt)
    text = response.text.strip()
    start = text.find('[')
    end = text.rfind(']') + 1
    if start != -1 and end != 0:
        doc_data = json.loads(text[start:end])
    else:
        raise ValueError("JSON topilmadi")
    
    doc = Document()
    doc.add_heading(topic, 0)
    for d in doc_data:
        doc.add_heading(d.get('heading', ''), level=1)
        doc.add_paragraph(d.get('text', ''))
    path = os.path.join(DOWNLOAD_DIR, f"doc_{int(time.time())}.docx")
    doc.save(path)
    return path

def search_youtube(query):
    """40 ta natija qidirish (4 sahifa uchun)"""
    ydl_opts = {'format': 'bestaudio/best', 'quiet': True, 'extract_flat': True, 'default_search': 'ytsearch', 'noplaylist': True}
    with yt_dlp.YoutubeDL(ydl_opts) as ydl:
        try:
            # 40 ta natija (har sahifada 10 tadan)
            return [e for e in ydl.extract_info(f"ytsearch40:{query}", download=False).get('entries', []) if e]
        except: return []

def download_audio(video_id, chat_id):
    """FFmpeg yordamida eng sifatli MP3 yuklash"""
    path_template = os.path.join(DOWNLOAD_DIR, f"{chat_id}_{video_id}")
    ydl_opts = {
        'format': 'bestaudio/best',
        'outtmpl': path_template,
        'quiet': True,
        'postprocessors': [{
            'key': 'FFmpegExtractAudio',
            'preferredcodec': 'mp3',
            'preferredquality': '192',
        }],
    }
    with yt_dlp.YoutubeDL(ydl_opts) as ydl:
        try:
            info = ydl.extract_info(f"https://www.youtube.com/watch?v={video_id}", download=True)
            # yt-dlp mp3 qo'shimchasini o'zi qo'shadi
            final_path = f"{path_template}.mp3"
            if not os.path.exists(final_path):
                # Agar kutilmaganda mp3 bo'lmasa, bor faylni topamiz
                possible_files = [f for f in os.listdir(DOWNLOAD_DIR) if f.startswith(f"{chat_id}_{video_id}")]
                if possible_files:
                    final_path = os.path.join(DOWNLOAD_DIR, possible_files[0])
            return final_path, info
        except Exception as e:
            print(f"❌ Yuklab olishda xato: {e}")
            raise e

def get_main_menu():
    """Asosiy menyu tugmalarini yaratish"""
    markup = types.ReplyKeyboardMarkup(resize_keyboard=True, row_width=2)
    item1 = types.KeyboardButton("🎵 Musiqa qidirish")
    item2 = types.KeyboardButton("🤖 AI Savol-javob")
    item3 = types.KeyboardButton("📊 Prezentatsiya yaratish")
    item4 = types.KeyboardButton("📄 Word hujjat tayyorlash")
    item5 = types.KeyboardButton("🏠 Asosiy menyu")
    markup.add(item1, item2, item3, item4)
    markup.row(item5)
    return markup

def get_menu_markup(entries, page=0):
    markup = types.InlineKeyboardMarkup(row_width=5)
    start = page * 10
    current_list = entries[start:start+10]
    
    # Raqamli tugmalar (tanlash uchun)
    btns = []
    for i, _ in enumerate(current_list):
        btns.append(types.InlineKeyboardButton(str(i+1), callback_data=f"download_{current_list[i]['id']}"))
    markup.add(*btns)
    
    # Navigatsiya tugmalari (Sahifalar)
    nav = []
    if page > 0:
        nav.append(types.InlineKeyboardButton("⬅️ Oldingi", callback_data=f"nav_{page-1}"))
    
    nav.append(types.InlineKeyboardButton("❌ Yopish", callback_data="close"))
    
    if (page + 1) * 10 < len(entries):
        nav.append(types.InlineKeyboardButton("Keyingi ➡️", callback_data=f"nav_{page+1}"))
    
    markup.row(*nav)
    return markup

@bot.message_handler(commands=['start', 'help'])
def hi(m):
    user_states[m.chat.id] = None
    text = (
        "👋 <b>Assalomu alaykum! Toparchik AI botiga xush kelibsiz!</b>\n\n"
        "Men sizga musiqa qidirishda, AI yordamida savollarga javob berishda va "
        "professional hujjatlar yaratishda yordam beraman.\n\n"
        "👇 <b>Quyidagi imkoniyatlardan birini tanlang:</b>"
    )
    bot.send_message(m.chat.id, text, reply_markup=get_main_menu(), parse_mode="HTML")

@bot.message_handler(func=lambda m: True)
def handle_text(message):
    chat_id = message.chat.id
    query = message.text.strip()
    
    # Menyu tugmalarini tekshirish
    if query == "🎵 Musiqa qidirish":
        user_states[chat_id] = 'music'
        return bot.send_message(chat_id, "🎵 <b>Musiqa nomini yozing:</b>", parse_mode="HTML")
    elif query == "🤖 AI Savol-javob":
        user_states[chat_id] = 'ai'
        return bot.send_message(chat_id, "🤖 <b>Istalgan savolingizni bering:</b>", parse_mode="HTML")
    elif query == "📊 Prezentatsiya yaratish":
        user_states[chat_id] = 'pptx'
        return bot.send_message(chat_id, "📊 <b>Prezentatsiya mavzusini yoki rejasini yozing:</b>", parse_mode="HTML")
    elif query == "📄 Word hujjat tayyorlash":
        user_states[chat_id] = 'docx'
        return bot.send_message(chat_id, "📄 <b>Hujjat (insho, referat) mavzusini yozing:</b>", parse_mode="HTML")
    elif query == "🏠 Asosiy menyu":
        user_states[chat_id] = None
        return bot.send_message(chat_id, "🏡 <b>Asosiy menyuga qaytdik. Bironta xizmatni tanlang:</b>", reply_markup=get_main_menu(), parse_mode="HTML")

    # Joriy holatni (mode) aniqlash
    mode = user_states.get(chat_id)
    
    if mode == 'pptx' or mode == 'docx':
        waiting_text = "📊 <b>Prezentatsiya slaydlari yaratilmoqda...</b>" if mode == 'pptx' else "✍️ <b>Referat/Hujjat tayyorlanmoqda...</b>"
        wait = bot.reply_to(message, f"⏳ {waiting_text}", parse_mode="HTML")
        try:
            path = generate_presentation(query) if mode == 'pptx' else generate_word_doc(query)
            with open(path, 'rb') as f:
                bot.send_document(message.chat.id, f)
            os.remove(path)
        except Exception as e:
            if "429" in str(e) or "quota" in str(e).lower():
                bot.send_message(message.chat.id, "😔 <b>Limit tugadi!</b> 1 daqiqa kuting.")
            else: bot.send_message(message.chat.id, "😔 Xatolik yuz berdi.")
        finally: 
            try: bot.delete_message(message.chat.id, wait.message_id)
            except: pass
        return

    if mode == 'ai':
        wait = bot.reply_to(message, "🤖 <b>Gemini o'ylamoqda...</b>", parse_mode="HTML")
        try:
            bot.send_message(message.chat.id, ask_gemini(query), parse_mode="HTML")
        except Exception as e:
            if "429" in str(e) or "quota" in str(e).lower():
                bot.send_message(message.chat.id, "😔 <b>AI Limiti tugadi!</b>")
            else: bot.send_message(message.chat.id, "😔 Xatolik.")
        finally: 
            try: bot.delete_message(message.chat.id, wait.message_id)
            except: pass
        return

    if mode == 'music' or not mode:
        wait = bot.reply_to(message, f"🔍 <b>Qidirilmoqda:</b> {query}...", parse_mode="HTML")
        results = search_youtube(query)
        try: bot.delete_message(message.chat.id, wait.message_id)
        except: pass
        if not results: return bot.send_message(message.chat.id, "😕 Topilmadi.")
        
        # Natijalarni saqlash
        user_searches[message.chat.id] = {"results": results, "query": query}
        save_user_searches() # Faylga yozish
        
        resp = f"🔍 <b>Natijalar:</b> {query}\n\n"
        for i, e in enumerate(results[:10], 1):
            resp += f"{i}. {e.get('title', '')[:40]} [{format_duration(e.get('duration'))}]\n"
        bot.send_message(message.chat.id, resp, reply_markup=get_menu_markup(results), parse_mode="HTML")

@bot.callback_query_handler(func=lambda c: True)
def calls(call):
    if call.data == "close": bot.delete_message(call.message.chat.id, call.message.message_id)
    elif call.data.startswith("download_"):
        vid = call.data.split("_")[1]
        bot.answer_callback_query(call.id, "📥 Musiqa tayyorlanmoqda...")
        wait = bot.send_message(call.message.chat.id, "🎵 <b>Fayl yuklab olinmoqda va konvertatsiya qilinmoqda...</b>", parse_mode="HTML")
        try:
            path, info = download_audio(vid, call.message.chat.id)
            with open(path, 'rb') as f:
                bot.send_audio(call.message.chat.id, f, title=info.get('title'))
            os.remove(path)
        except: bot.send_message(call.message.chat.id, "❌ Xato.")
        finally:
            try: bot.delete_message(call.message.chat.id, wait.message_id)
            except: pass
    elif call.data.startswith("nav_"):
        page = int(call.data.split("_")[1])
        ds = user_searches.get(call.message.chat.id)
        if not ds: return bot.answer_callback_query(call.id, "😕 Qidiruv muddati o'tgan. Iltimos, qaytadan qidiring.")
        resp = f"🔍 <b>Natijalar:</b> {ds['query']} ({page+1}-bet)\n\n"
        for i, e in enumerate(ds["results"][page*10:(page+1)*10], 1):
            resp += f"{page*10+i}. {e.get('title', '')[:40]}\n"
        bot.edit_message_text(resp, call.message.chat.id, call.message.message_id, reply_markup=get_menu_markup(ds["results"], page), parse_mode="HTML")

if __name__ == "__main__":
    log("🤖 Bot polling rejimida ishga tushmoqda...")
    while True:
        try:
            bot.infinity_polling(timeout=90, long_polling_timeout=40)
        except Exception as e:
            log(f"⚠️ Xato: {e}")
            time.sleep(10)